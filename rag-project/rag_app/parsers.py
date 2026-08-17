from __future__ import annotations

import csv
import io
from pathlib import Path

from .chunking import TextUnit


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".html", ".htm"}


def _read_text(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "big5"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_units(path: Path, original_name: str) -> list[TextUnit]:
    ext = Path(original_name).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{ext or '无扩展名'}")

    if ext in {".txt", ".md"}:
        return [TextUnit(_read_text(path), "全文")]
    if ext in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        text = BeautifulSoup(_read_text(path), "html.parser").get_text("\n", strip=True)
        return [TextUnit(text, "网页正文")]
    if ext == ".pdf":
        from pypdf import PdfReader

        units: list[TextUnit] = []
        for number, page in enumerate(PdfReader(str(path)).pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                units.append(TextUnit(text, f"第 {number} 页"))
        return units
    if ext == ".docx":
        from docx import Document

        document = Document(str(path))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return [TextUnit("\n".join(parts), "文档正文")]
    if ext == ".pptx":
        from pptx import Presentation

        units = []
        for number, slide in enumerate(Presentation(str(path)).slides, 1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            if text:
                units.append(TextUnit(text, f"第 {number} 页幻灯片"))
        return units
    if ext == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        units = []
        for sheet in workbook.worksheets:
            rows = [" | ".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True)]
            text = "\n".join(row for row in rows if row.strip(" |"))
            if text:
                units.append(TextUnit(text, f"工作表 {sheet.title}"))
        return units
    if ext == ".csv":
        rows = [" | ".join(row) for row in csv.reader(io.StringIO(_read_text(path)))]
        return [TextUnit("\n".join(rows), "CSV 表格")]
    raise ValueError(f"无法解析文件：{original_name}")
